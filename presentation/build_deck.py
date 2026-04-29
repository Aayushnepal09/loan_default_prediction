"""
Build Phase 4 presentation deck (Phase4_presentation.pptx).

Design system:
  - 16:9 widescreen, 13.333" x 7.5"
  - Palette:
      Navy primary  #0E3D70   - dominates titles/headers
      Deep navy bg  #11294D   - title + closing slides
      Light grey bg #F4F6FA   - subtle backgrounds and chip strips
      Amber accent  #E89E2C   - callout numbers and "PHASE N" chips
      Risk red      #D73027   - matches figures
      Safe green    #1A9850   - matches figures
      Body text     #1A2440
  - Visual motif: a small "PHASE N" chip top-right on every phase recap slide
                  + thin colored bar across the slide top
  - Typography: Calibri throughout; Calibri Light for body, Calibri Bold for headers.
                Big stat callouts in 60-72pt.

Real metrics from generate_figures.py output:
  Test set 314,212 loans (2017), 21.0% default rate
  AUC-ROC 0.7260, AUC-PR 0.4103, KS 0.3279
  Threshold 0.4951 -> Acc 65.9%, Prec 34.2%, Recall 67.3%, F1 0.453
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


ROOT = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = ROOT / "Phase4_presentation.pptx"


# Palette
NAVY        = RGBColor(0x0E, 0x3D, 0x70)
NAVY_DEEP   = RGBColor(0x11, 0x29, 0x4D)
GREY_BG     = RGBColor(0xF4, 0xF6, 0xFA)
AMBER       = RGBColor(0xE8, 0x9E, 0x2C)
RED         = RGBColor(0xD7, 0x30, 0x27)
GREEN       = RGBColor(0x1A, 0x98, 0x50)
BODY        = RGBColor(0x1A, 0x24, 0x40)
MUTED       = RGBColor(0x55, 0x60, 0x75)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LINE_GREY   = RGBColor(0xCC, 0xD4, 0xE0)


SLIDE_W = 13.333
SLIDE_H = 7.5


def hex_for(rgb: RGBColor) -> str:
    return "%02X%02X%02X" % (rgb[0], rgb[1], rgb[2])


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    if fill_rgb is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        if line_width is not None:
            shape.line.width = Pt(line_width)
    return shape


def add_round_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=None,
                    radius=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    # set corner radius via adjustments
    shape.adjustments[0] = radius
    if fill_rgb is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        if line_width is not None:
            shape.line.width = Pt(line_width)
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=BODY,
              face="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = face
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_phase_chip(slide, phase_num: int, x=12.1, y=0.32):
    """Small amber pill 'PHASE N' top-right."""
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(0.95), Inches(0.32))
    pill.shadow.inherit = False
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = AMBER
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"PHASE {phase_num}"
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_top_bar(slide, color=NAVY, h=0.12):
    add_rect(slide, 0, 0, SLIDE_W, h, fill_rgb=color)


def slide_blank(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def add_footer(slide, page_num: int, total: int):
    add_text(slide, 0.5, 7.16, 6, 0.25,
             "EAS 587 Phase 4  -  Lending Club Loan Default Prediction",
             size=9, color=MUTED)
    add_text(slide, SLIDE_W - 1.5, 7.16, 1.0, 0.25,
             f"{page_num} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ---- slide builders --------------------------------------------------------


def build_title(prs):
    s = slide_blank(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY_DEEP

    # left amber bar
    add_rect(s, 0, 0, 0.18, SLIDE_H, fill_rgb=AMBER)

    add_text(s, 0.9, 1.5, 12, 0.5,
             "EAS 587  -  Phase 4  -  Spring 2026",
             size=18, color=AMBER, bold=True, face="Calibri")

    add_text(s, 0.9, 2.1, 12, 1.6,
             "Lending Club\nLoan Default Prediction",
             size=58, color=WHITE, bold=True, face="Calibri")

    add_text(s, 0.9, 4.4, 12, 1.0,
             "Same use case, four implementations:\npandas  ->  scikit-learn  ->  Spark + MCP  ->  conversational UI",
             size=20, color=RGBColor(0xCA, 0xDC, 0xFC), italic=True)

    add_text(s, 0.9, 6.2, 12, 0.35,
             "Team  -  404 Team Not Found",
             size=13, color=AMBER, bold=True)
    add_text(s, 0.9, 6.55, 12, 0.4,
             "Aayush Nepal  |  Junwei Zhang  |  Lusi Zhang",
             size=14, color=WHITE)
    add_text(s, 0.9, 6.8, 12, 0.3,
             "Final presentations: week of May 4, 2026",
             size=12, color=RGBColor(0x9F, 0xB3, 0xD1), italic=True)


def build_use_case(prs):
    s = slide_blank(prs)
    add_top_bar(s)

    add_text(s, 0.6, 0.45, 12, 0.7, "The use case", size=34, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 12, 0.4,
             "A peer-to-peer lender's question, restated as a binary classification problem",
             size=14, color=MUTED, italic=True)

    # Left side: the question, framed
    add_round_rect(s, 0.6, 1.85, 6.2, 4.6, fill_rgb=GREY_BG, radius=0.04)
    add_text(s, 0.95, 2.05, 5.5, 0.5,
             "The lender's question", size=14, bold=True, color=NAVY)
    add_text(s, 0.95, 2.55, 5.5, 1.4,
             "\"Given a loan application, what is the probability\nthis borrower will charge off rather than fully repay?\"",
             size=20, italic=True, color=BODY)

    add_text(s, 0.95, 4.1, 5.5, 0.4,
             "Decisions the model supports", size=12, bold=True, color=NAVY)
    bullets = [
        "Approve / decline at origination",
        "Risk-based pricing (interest rate)",
        "Portfolio-level default rate forecast",
        "Manual-review queue for borderline loans",
    ]
    for i, b in enumerate(bullets):
        add_text(s, 1.15, 4.5 + i * 0.36, 5.3, 0.36,
                 "-  " + b, size=14, color=BODY)

    # Right side: stat callouts (3 stacked)
    stats = [
        ("314,212", "completed loans 2017 (test set)",     NAVY),
        ("21.0%",   "actual charge-off rate",              RED),
        ("~4 to 1", "fully paid : default class imbalance",  AMBER),
    ]
    base_y = 1.85
    for i, (big, small, color) in enumerate(stats):
        y = base_y + i * 1.55
        add_round_rect(s, 7.2, y, 5.5, 1.4, fill_rgb=WHITE,
                       line_rgb=LINE_GREY, line_width=0.75, radius=0.05)
        add_text(s, 7.4, y + 0.18, 5.1, 0.7, big,
                 size=42, bold=True, color=color)
        add_text(s, 7.4, y + 0.85, 5.1, 0.45, small,
                 size=13, color=MUTED)

    add_footer(s, 2, 13)


def build_evolution(prs):
    s = slide_blank(prs)
    add_top_bar(s)

    add_text(s, 0.6, 0.45, 12, 0.7, "Project evolution", size=34, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 12, 0.4,
             "One question, four artifacts of increasing scale and reach",
             size=14, color=MUTED, italic=True)

    s.shapes.add_picture(str(FIG / "architecture_diagram.png"),
                          Inches(0.5), Inches(1.9), width=Inches(12.3))

    add_text(s, 0.6, 6.5, 12.2, 0.5,
             "Phases 1+2 received perfect scores; we retained the time-based split, leakage discipline, and 8-stage pipeline and built on top.",
             size=13, italic=True, color=MUTED)

    add_footer(s, 3, 13)


def build_phase1(prs):
    s = slide_blank(prs)
    add_top_bar(s)
    add_phase_chip(s, 1)

    add_text(s, 0.6, 0.45, 11, 0.7, "Data foundations", size=34, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 11, 0.4,
             "Filter, define the target, and split with leakage discipline",
             size=14, color=MUTED, italic=True)

    # 3 columns: Filter, Target, Split
    cols = [
        ("FILTER  -  2014 to 2017",
         [
             "Pre-2014 issuance volume too sparse to be representative.",
             "Post-2017 loans not matured: charge-off labels unreliable.",
             "~1.2M completed loans after filter, sorted chronologically.",
         ]),
        ("TARGET  -  charged_off",
         [
             "Keep only loans with terminal status (Fully Paid / Charged Off).",
             "Class balance: ~80% paid / 20% default.",
             "No SMOTE or class re-weighting -- model handles imbalance natively.",
         ]),
        ("SPLIT  -  time-based",
         [
             "Train: first 80% of 2014 to 2016 by issue_d.",
             "Val:   last 20% of 2014 to 2016.",
             "Test:  ALL of 2017 (held out, ~314k loans).",
             "Random splits would leak the macro environment.",
         ]),
    ]
    base_x = 0.6
    col_w = 4.05
    gap = 0.1
    for i, (title, body_lines) in enumerate(cols):
        x = base_x + i * (col_w + gap)
        add_round_rect(s, x, 1.85, col_w, 3.4, fill_rgb=GREY_BG, radius=0.04)
        add_rect(s, x, 1.85, 0.07, 3.4, fill_rgb=NAVY)  # left accent
        add_text(s, x + 0.25, 2.0, col_w - 0.4, 0.4, title,
                 size=13, bold=True, color=NAVY)
        for j, line in enumerate(body_lines):
            add_text(s, x + 0.25, 2.5 + j * 0.6, col_w - 0.4, 0.6,
                     "-  " + line, size=12, color=BODY)

    # leakage callout strip
    add_round_rect(s, 0.6, 5.5, 12.13, 1.45, fill_rgb=WHITE,
                   line_rgb=AMBER, line_width=1.5, radius=0.05)
    add_text(s, 0.85, 5.7, 11.9, 0.4,
             "Leakage discipline -- our most defensible Phase 1 decision",
             size=14, bold=True, color=NAVY)
    add_text(s, 0.85, 6.15, 11.9, 0.7,
             "We dropped 21 post-loan columns that only get values AFTER the outcome is known: total_pymnt*, recoveries, last_pymnt_*, last_fico_range_*, hardship_flag, and full settlement details. Most public Lending Club notebooks accidentally train on these.",
             size=12, color=BODY)

    add_footer(s, 4, 13)


def build_phase2_pipeline(prs):
    s = slide_blank(prs)
    add_top_bar(s)
    add_phase_chip(s, 2)

    add_text(s, 0.6, 0.45, 11, 0.7,
             "8-stage preprocessing pipeline", size=32, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 11, 0.4,
             "Custom scikit-learn transformers, fit on train only, applied identically to val + test",
             size=14, color=MUTED, italic=True)

    s.shapes.add_picture(str(FIG / "pipeline_diagram.png"),
                          Inches(0.4), Inches(1.85), width=Inches(12.5))

    # Two-column "why each stage" callout
    add_round_rect(s, 0.6, 5.5, 6.05, 1.45, fill_rgb=GREY_BG, radius=0.04)
    add_text(s, 0.85, 5.65, 5.7, 0.4,
             "Engineered features", size=13, bold=True, color=NAVY)
    add_text(s, 0.85, 6.05, 5.7, 0.85,
             "loan_to_income, monthly_payment_to_income, delinq_rate,\ncredit_history_months, issue_year, issue_month",
             size=11.5, color=BODY)

    add_round_rect(s, 6.85, 5.5, 6.05, 1.45, fill_rgb=GREY_BG, radius=0.04)
    add_text(s, 7.1, 5.65, 5.7, 0.4,
             "Macro overlay", size=13, bold=True, color=NAVY)
    add_text(s, 7.1, 6.05, 5.7, 0.85,
             "FRED UNRATE merged on issue (year, month).\nSecondary data source flagged in the rubric.",
             size=11.5, color=BODY)

    add_footer(s, 5, 13)


def build_phase2_modelsel(prs):
    s = slide_blank(prs)
    add_top_bar(s)
    add_phase_chip(s, 2)

    add_text(s, 0.6, 0.45, 11, 0.7,
             "Model bake-off", size=34, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 11, 0.4,
             "Four algorithms, 50 Optuna trials each, MLflow tracked, validated on the 2014-16 hold-out",
             size=14, color=MUTED, italic=True)

    s.shapes.add_picture(str(FIG / "model_comparison.png"),
                          Inches(0.5), Inches(1.85), width=Inches(7.6))

    # Right side: winner card
    add_round_rect(s, 8.4, 1.85, 4.5, 4.4, fill_rgb=NAVY, radius=0.05)
    add_text(s, 8.65, 2.05, 4.0, 0.4,
             "WINNER", size=12, bold=True, color=AMBER)
    add_text(s, 8.65, 2.5, 4.0, 0.6,
             "XGBoost", size=34, bold=True, color=WHITE)
    add_text(s, 8.65, 3.2, 4.0, 0.4,
             "Why it won", size=12, bold=True,
             color=RGBColor(0xCA, 0xDC, 0xFC))
    why = [
        "Highest validation AUC-ROC (0.7238).",
        "Highest KS statistic (0.325).",
        "5x faster than HGB at fit time.",
        "Native pred_contribs (SHAP) for explainability.",
    ]
    for i, line in enumerate(why):
        add_text(s, 8.85, 3.6 + i * 0.45, 3.8, 0.45,
                 "-  " + line, size=12, color=WHITE)

    # bottom strip with method
    add_round_rect(s, 0.6, 6.5, 12.13, 0.55, fill_rgb=GREY_BG, radius=0.05)
    add_text(s, 0.85, 6.58, 12, 0.4,
             "Tuning: Optuna TPE, 50 trials per model, optimizing val AUC-ROC. All 200+ trials logged to MLflow at models/mlruns/.",
             size=12, color=BODY)

    add_footer(s, 6, 13)


def build_phase2_results(prs):
    s = slide_blank(prs)
    add_top_bar(s)
    add_phase_chip(s, 2)

    add_text(s, 0.6, 0.45, 11, 0.7,
             "Hold-out test results", size=34, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 11, 0.4,
             "All 314,212 loans issued in 2017 -- never seen during training or tuning",
             size=14, color=MUTED, italic=True)

    # 4 stat callouts in a row
    stats = [
        ("0.726", "AUC-ROC",  NAVY),
        ("0.410", "AUC-PR",   NAVY),
        ("0.328", "KS",       AMBER),
        ("67.3%", "Recall on defaults\n(at Youden threshold)", RED),
    ]
    for i, (big, small, color) in enumerate(stats):
        x = 0.6 + i * 3.13
        add_round_rect(s, x, 1.8, 2.95, 1.3, fill_rgb=WHITE,
                       line_rgb=LINE_GREY, line_width=0.75, radius=0.05)
        add_text(s, x + 0.15, 1.92, 2.65, 0.7, big,
                 size=38, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.15, 2.6, 2.65, 0.55, small,
                 size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

    # ROC curve and confusion matrix
    s.shapes.add_picture(str(FIG / "roc_curve.png"),
                          Inches(0.6), Inches(3.3), height=Inches(3.6))
    s.shapes.add_picture(str(FIG / "confusion_matrix.png"),
                          Inches(8.0), Inches(3.3), height=Inches(3.6))

    add_footer(s, 7, 13)


def build_phase3_spark(prs):
    s = slide_blank(prs)
    add_top_bar(s)
    add_phase_chip(s, 3)

    add_text(s, 0.6, 0.45, 11, 0.7,
             "Scaling on Spark + Delta Lake", size=32, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 11, 0.4,
             "Same data, same target, rebuilt on Databricks using the medallion architecture",
             size=14, color=MUTED, italic=True)

    # 3 medallion cards
    layers = [
        ("BRONZE",
         "bronze_loans",
         "Kaggle download landed as-is.\nFiltered 2014-2017.\nWritten to Unity Catalog volume.",
         RGBColor(0xCD, 0x7F, 0x32)),
        ("SILVER",
         "silver_loans",
         "Cleaned and type-fixed.\nLeakage columns dropped.\nSchema enforced.",
         RGBColor(0xA0, 0xA8, 0xB0)),
        ("GOLD",
         "gold_loans_train / val / test",
         "Time-based split.\nStratified samples for train/val.\nFull 2017 test set retained.",
         RGBColor(0xC9, 0xA1, 0x3B)),
    ]
    for i, (label, table, body, accent) in enumerate(layers):
        x = 0.6 + i * 4.13
        add_round_rect(s, x, 1.85, 4.0, 3.3, fill_rgb=GREY_BG, radius=0.05)
        add_rect(s, x, 1.85, 4.0, 0.4, fill_rgb=accent)
        add_text(s, x + 0.2, 1.92, 3.6, 0.3, label,
                 size=13, bold=True, color=WHITE)
        add_text(s, x + 0.2, 2.4, 3.6, 0.4, table,
                 size=14, bold=True, color=NAVY, face="Consolas")
        add_text(s, x + 0.2, 2.95, 3.6, 2.0, body,
                 size=12, color=BODY)

    # Why Spark callout
    add_round_rect(s, 0.6, 5.45, 12.13, 1.5, fill_rgb=WHITE,
                   line_rgb=NAVY, line_width=1.0, radius=0.05)
    add_text(s, 0.85, 5.6, 11.9, 0.4,
             "Why rebuild on Spark", size=14, bold=True, color=NAVY)
    add_text(s, 0.85, 6.05, 11.9, 0.85,
             "Phase 1+2 ran in pandas on a laptop and would not survive the full 2.9M-row Lending Club archive. Bronze->Silver->Gold gives us reproducible, versioned tables and lets us re-run the whole pipeline on the cluster without touching local CSVs.",
             size=12, color=BODY)

    add_footer(s, 8, 13)


def build_phase3_macro_mcp(prs):
    s = slide_blank(prs)
    add_top_bar(s)
    add_phase_chip(s, 3)

    add_text(s, 0.6, 0.45, 11, 0.7,
             "Macro overlay + MCP deployment", size=30, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 11, 0.4,
             "Two ideas borrowed from outside the loan dataset itself",
             size=14, color=MUTED, italic=True)

    # Left: FRED macro
    add_round_rect(s, 0.6, 1.85, 6.05, 5.0, fill_rgb=GREY_BG, radius=0.04)
    add_rect(s, 0.6, 1.85, 0.07, 5.0, fill_rgb=AMBER)
    add_text(s, 0.85, 2.0, 5.7, 0.4,
             "Secondary data source: FRED", size=14, bold=True, color=NAVY)
    add_text(s, 0.85, 2.45, 5.7, 0.4,
             "U.S. monthly unemployment (UNRATE)", size=12, italic=True, color=MUTED)
    fred_lines = [
        "Pulled at fit-time via pandas-datareader.",
        "Joined on (issue_year, issue_month) per loan.",
        "Captures macro pressure that borrower features cannot.",
        "Falls back to 5.0% if FRED is unreachable -- prevents pipeline crash on offline runs.",
    ]
    for i, line in enumerate(fred_lines):
        add_text(s, 1.0, 3.0 + i * 0.7, 5.5, 0.7,
                 "-  " + line, size=12.5, color=BODY)
    add_text(s, 0.85, 6.15, 5.7, 0.65,
             "Why it matters: 2014-2017 spans a falling-unemployment regime; the model learns that loans issued during weaker labor markets default more often.",
             size=11, italic=True, color=MUTED)

    # Right: MCP server
    add_round_rect(s, 6.85, 1.85, 6.05, 5.0, fill_rgb=NAVY, radius=0.04)
    add_text(s, 7.1, 2.0, 5.7, 0.4,
             "MCP deployment server", size=14, bold=True,
             color=AMBER)
    add_text(s, 7.1, 2.45, 5.7, 0.4,
             "FastMCP + Claude Desktop", size=12, italic=True,
             color=RGBColor(0xCA, 0xDC, 0xFC))
    mcp_lines = [
        "Loads best_model.pkl and preprocessing_pipeline.pkl on startup.",
        "Exposes predict_loan_default tool over stdio MCP.",
        "Validates inputs, fills 24 secondary defaults, runs the same transform path as training.",
        "Returns probability + risk tier + decision recommendation.",
        "Demonstrated live in the Phase 4 demo.",
    ]
    for i, line in enumerate(mcp_lines):
        add_text(s, 7.25, 3.0 + i * 0.6, 5.55, 0.6,
                 "-  " + line, size=12, color=WHITE)

    add_footer(s, 9, 13)


def build_demo(prs):
    s = slide_blank(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY_DEEP

    add_phase_chip(s, 4)

    add_text(s, 0.6, 0.7, 12, 0.7,
             "LIVE DEMO", size=20, bold=True, color=AMBER)
    add_text(s, 0.6, 1.2, 12, 1.4,
             "Streamlit predictor\n+ MCP cherry-on-top",
             size=46, bold=True, color=WHITE)

    # demo storyboard - 4 numbered steps
    steps = [
        ("1", "Pick a preset",
         "Sidebar: Safe / Borderline / Risky profiles for instant fill."),
        ("2", "Predict",
         "Plotly gauge + 5-tier risk badge + decision recommendation."),
        ("3", "Why this prediction",
         "Top 8 SHAP contributions from XGBoost's native pred_contribs."),
        ("4", "Same model, conversational",
         "Switch to Claude Desktop -- describe a loan in plain English, MCP runs the same pipeline."),
    ]
    base_y = 3.4
    for i, (num, title, body) in enumerate(steps):
        y = base_y + (i // 2) * 1.8
        x = 0.6 + (i % 2) * 6.3
        add_round_rect(s, x, y, 6.05, 1.55, fill_rgb=RGBColor(0x1B, 0x3A, 0x65),
                       radius=0.04)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x + 0.25), Inches(y + 0.25),
                                    Inches(0.7), Inches(0.7))
        circ.shadow.inherit = False
        circ.fill.solid(); circ.fill.fore_color.rgb = AMBER
        circ.line.fill.background()
        ctf = circ.text_frame
        ctf.margin_left = 0; ctf.margin_right = 0
        ctf.margin_top = 0; ctf.margin_bottom = 0
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run(); crun.text = num
        crun.font.name = "Calibri"; crun.font.size = Pt(22); crun.font.bold = True
        crun.font.color.rgb = NAVY_DEEP

        add_text(s, x + 1.1, y + 0.22, 4.7, 0.45, title,
                 size=16, bold=True, color=WHITE)
        add_text(s, x + 1.1, y + 0.7, 4.7, 0.85, body,
                 size=11.5, color=RGBColor(0xCA, 0xDC, 0xFC))

    add_text(s, 0.6, 7.05, 12, 0.3,
             "If the live demo fails: Streamlit screenshots in the appendix; MCP can be invoked as a backup prompt.",
             size=10, italic=True, color=RGBColor(0x7A, 0x8E, 0xAE))


def build_insights(prs):
    s = slide_blank(prs)
    add_top_bar(s)

    add_text(s, 0.6, 0.45, 12, 0.7, "Insights", size=34, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 12, 0.4,
             "What the model learned -- and what that says about the lending market",
             size=14, color=MUTED, italic=True)

    # Left: feature importance image
    s.shapes.add_picture(str(FIG / "feature_importance.png"),
                          Inches(0.5), Inches(1.85), height=Inches(5.0))

    # Right: insight cards
    insights = [
        ("Sub-grade dominates",
         "sub_grade carries ~3x more gain than the next feature. Lending Club's own grade letter is already a near-sufficient summary of borrower risk -- our pipeline mostly refines it."),
        ("Term length is a hidden risk lever",
         "Longer 60-month loans default substantially more often than 36-month loans at the same grade. The model picks this up without being told."),
        ("Macro context matters",
         "unemployment_rate (FRED) consistently appears in the top contributions. Removing it dropped val AUC by ~0.5 points -- secondary data earns its keep."),
        ("Renting > owning, all else equal",
         "home_ownership = RENT pushes toward default; MORTGAGE pulls away. Likely a wealth proxy."),
    ]
    base_y = 1.85
    card_h = 1.2
    for i, (t, body) in enumerate(insights):
        y = base_y + i * (card_h + 0.08)
        add_round_rect(s, 7.5, y, 5.4, card_h, fill_rgb=GREY_BG, radius=0.04)
        add_rect(s, 7.5, y, 0.07, card_h, fill_rgb=AMBER)
        add_text(s, 7.7, y + 0.1, 5.1, 0.35, t,
                 size=13, bold=True, color=NAVY)
        add_text(s, 7.7, y + 0.45, 5.1, card_h - 0.5, body,
                 size=11, color=BODY)

    add_footer(s, 11, 13)


def build_limitations(prs):
    s = slide_blank(prs)
    add_top_bar(s)

    add_text(s, 0.6, 0.45, 12, 0.7,
             "Limitations + what's next", size=32, bold=True, color=NAVY)
    add_text(s, 0.6, 1.1, 12, 0.4,
             "Honest about what we did NOT solve, and where this work goes",
             size=14, color=MUTED, italic=True)

    # Two columns
    add_round_rect(s, 0.6, 1.85, 6.05, 5.0, fill_rgb=GREY_BG, radius=0.04)
    add_rect(s, 0.6, 1.85, 0.07, 5.0, fill_rgb=RED)
    add_text(s, 0.9, 2.0, 5.6, 0.4, "Limitations",
             size=18, bold=True, color=NAVY)
    lims = [
        ("Class imbalance left as-is",
         "We did not resample. Any deployment with a different operating cost-ratio needs threshold re-tuning."),
        ("Concept drift after 2017",
         "Training cohort spans falling unemployment. Model is not validated on COVID-era loan books."),
        ("No fairness audit",
         "Lending data has well-documented protected-class disparities. Disparate-impact testing was out of scope."),
        ("FICO-only credit signal",
         "We use fico_range_low; richer bureau features would likely add lift but were absent from this dataset."),
    ]
    for i, (t, body) in enumerate(lims):
        y = 2.5 + i * 1.05
        add_text(s, 0.95, y, 5.55, 0.35, "-  " + t,
                 size=13, bold=True, color=BODY)
        add_text(s, 1.2, y + 0.4, 5.3, 0.6, body,
                 size=11, color=MUTED)

    add_round_rect(s, 6.85, 1.85, 6.05, 5.0, fill_rgb=GREY_BG, radius=0.04)
    add_rect(s, 6.85, 1.85, 0.07, 5.0, fill_rgb=GREEN)
    add_text(s, 7.15, 2.0, 5.6, 0.4, "What's next",
             size=18, bold=True, color=NAVY)
    nexts = [
        ("Real-time scoring API",
         "Wrap the same pipeline in a FastAPI service; the MCP server already shows the loading pattern."),
        ("Fairness + adverse-action layer",
         "Disparate-impact metrics by protected group + per-prediction reason codes for compliant denial letters."),
        ("Retraining cadence",
         "Quarterly re-fit on the latest closed cohort; alarm on AUC drift > 0.02."),
        ("More macro signals",
         "FRED unemployment is one series; rates, housing, and consumer sentiment likely add lift."),
    ]
    for i, (t, body) in enumerate(nexts):
        y = 2.5 + i * 1.05
        add_text(s, 7.2, y, 5.55, 0.35, "-  " + t,
                 size=13, bold=True, color=BODY)
        add_text(s, 7.45, y + 0.4, 5.3, 0.6, body,
                 size=11, color=MUTED)

    add_footer(s, 12, 13)


def build_close(prs):
    s = slide_blank(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY_DEEP

    add_rect(s, 0, 0, 0.18, SLIDE_H, fill_rgb=AMBER)

    add_text(s, 0.9, 1.6, 12, 1.3,
             "Thank you", size=58, bold=True, color=WHITE)
    add_text(s, 0.9, 2.95, 12, 0.7,
             "Questions?", size=30, color=AMBER)
    add_text(s, 0.9, 3.7, 12, 0.4,
             "404 Team Not Found  -  Aayush Nepal | Junwei Zhang | Lusi Zhang",
             size=13, color=WHITE, italic=True)

    # Three takeaway cards on dark
    takeaways = [
        ("WHAT WE BUILT",
         "End-to-end loan default predictor: pandas pipeline, Spark medallion, MCP server, Streamlit app."),
        ("HEADLINE METRIC",
         "AUC-ROC 0.726 on 314k held-out 2017 loans -- never seen during training or tuning."),
        ("WHY IT MATTERS",
         "The same model is reachable two ways: a visual UI for analysts, a conversational tool for everyone else."),
    ]
    for i, (t, body) in enumerate(takeaways):
        x = 0.9 + i * 4.15
        add_round_rect(s, x, 4.5, 4.0, 2.0, fill_rgb=RGBColor(0x1B, 0x3A, 0x65),
                       radius=0.04)
        add_text(s, x + 0.25, 4.65, 3.7, 0.4, t,
                 size=12, bold=True, color=AMBER)
        add_text(s, x + 0.25, 5.1, 3.7, 1.3, body,
                 size=12, color=WHITE)

    add_text(s, 0.9, 7.0, 12, 0.3,
             "Repository: github.com/Aayushnepal09/Eas587_project    |    EAS 587 Spring 2026",
             size=11, color=RGBColor(0x9F, 0xB3, 0xD1), italic=True)


def main():
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title(prs)
    build_use_case(prs)
    build_evolution(prs)
    build_phase1(prs)
    build_phase2_pipeline(prs)
    build_phase2_modelsel(prs)
    build_phase2_results(prs)
    build_phase3_spark(prs)
    build_phase3_macro_mcp(prs)
    build_demo(prs)
    build_insights(prs)
    build_limitations(prs)
    build_close(prs)

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
